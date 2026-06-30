"""
Generate a 7-slide PowerPoint presentation from the research paper.
Slide 1: Title / Introduction
Slide 2: Research Background & Related Work
Slide 3: Dataset & Data Pipeline
Slide 4: Methodology
Slide 5: Results – EDA & Correlation
Slide 6: Results – Statistical Tests, ANOVA & Clustering
Slide 7: Conclusion & Future Work
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

BASE   = os.path.dirname(os.path.abspath(__file__))
FIGS   = os.path.join(BASE, "_paper_figures")
OUTPUT = os.path.join(BASE, "ClimateChange_Presentation.pptx")

# ── Colour palette ────────────────────────────────────────────────────────────
DARK_BLUE   = RGBColor(0x0D, 0x2B, 0x4E)   # slide background / headers
MID_BLUE    = RGBColor(0x1A, 0x57, 0x9E)   # accent bars
LIGHT_BLUE  = RGBColor(0xD6, 0xE8, 0xF7)   # subtle fill
GREEN       = RGBColor(0x27, 0xAE, 0x60)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE   = RGBColor(0xF2, 0xF6, 0xFC)
DARK_GREY   = RGBColor(0x2C, 0x3E, 0x50)
MID_GREY    = RGBColor(0x7F, 0x8C, 0x8D)

# ── Slide dimensions: widescreen 16:9 ────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank

# ═══════════════════════════════════════════════════════════════════════════
# HELPERS
# ═══════════════════════════════════════════════════════════════════════════

def add_rect(slide, x, y, w, h, fill, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    return shape

def add_text(slide, text, x, y, w, h,
             size=18, bold=False, italic=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox

def add_para(tf, text, size=14, bold=False, color=DARK_GREY,
             space_before=6, indent=False):
    p = tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    p.space_before = Pt(space_before)
    if indent:
        p.level = 1
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color

def slide_bg(slide, color=OFF_WHITE):
    bg = add_rect(slide, 0, 0, 13.33, 7.5, color)
    bg.zorder = 0

def header_bar(slide, title, subtitle=None):
    """Dark blue top bar with title."""
    add_rect(slide, 0, 0, 13.33, 1.3, DARK_BLUE)
    add_text(slide, title, 0.35, 0.1, 12.5, 0.75,
             size=30, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, 0.35, 0.78, 12.5, 0.42,
                 size=15, color=LIGHT_BLUE, align=PP_ALIGN.LEFT)
    # thin accent line
    add_rect(slide, 0, 1.3, 13.33, 0.06, MID_BLUE)

def bullet_box(slide, x, y, w, h, title, bullets,
               title_color=DARK_BLUE, bullet_size=14, bg=True):
    """White card with a bold title and bullet list."""
    if bg:
        card = add_rect(slide, x, y, w, h, WHITE)
        card.line.color.rgb = LIGHT_BLUE
        card.line.width = Pt(0.75)
    # title
    t = slide.shapes.add_textbox(Inches(x+0.12), Inches(y+0.1),
                                  Inches(w-0.2), Inches(0.38))
    tf = t.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(15)
    run.font.bold = True
    run.font.color.rgb = title_color
    # divider
    div = slide.shapes.add_shape(1,
              Inches(x+0.12), Inches(y+0.48),
              Inches(w-0.24), Inches(0.03))
    div.fill.solid(); div.fill.fore_color.rgb = MID_BLUE
    div.line.fill.background()
    # bullets
    b = slide.shapes.add_textbox(Inches(x+0.12), Inches(y+0.55),
                                  Inches(w-0.24), Inches(h-0.7))
    btf = b.text_frame
    btf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p2 = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
        p2.space_before = Pt(4)
        run2 = p2.add_run()
        run2.text = f"•  {bullet}"
        run2.font.size = Pt(bullet_size)
        run2.font.color.rgb = DARK_GREY

def img(slide, path, x, y, w):
    if os.path.exists(path):
        slide.shapes.add_picture(path, Inches(x), Inches(y), width=Inches(w))

def footer(slide, text="Data-Driven Analysis of Climate Change Indicators  |  CSUB & AVC"):
    add_rect(slide, 0, 7.18, 13.33, 0.32, DARK_BLUE)
    add_text(slide, text, 0.2, 7.2, 12.9, 0.28,
             size=9, color=MID_GREY, align=PP_ALIGN.LEFT)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 – INTRODUCTION / TITLE
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)

# Full dark-blue background
add_rect(sl, 0, 0, 13.33, 7.5, DARK_BLUE)
# Decorative right panel
add_rect(sl, 9.5, 0, 3.83, 7.5, MID_BLUE)
# Subtle diagonal accent
add_rect(sl, 8.9, 0, 0.12, 7.5, GREEN)

# Title
add_text(sl, "Data-Driven Analysis of", 0.5, 1.2, 9.0, 0.8,
         size=34, bold=True, color=WHITE)
add_text(sl, "Climate Change Indicators", 0.5, 1.9, 9.0, 0.8,
         size=34, bold=True, color=WHITE)

# Subtitle
add_text(sl, "A Multi-Layer Statistical Pipeline for Renewable Energy & CO₂ Analysis",
         0.5, 2.85, 8.8, 0.6, size=18, italic=True, color=LIGHT_BLUE)

# Divider
add_rect(sl, 0.5, 3.65, 5.5, 0.06, GREEN)

# Authors
add_text(sl, "Alex Vartanian\nCalifornia State University Bakersfield",
         0.5, 3.85, 5.5, 0.8, size=14, color=WHITE)
add_text(sl, "Ammar Abdelmoneam\nAntelope Valley College",
         0.5, 4.65, 5.5, 0.8, size=14, color=WHITE)

# Key facts on right panel
add_text(sl, "At a Glance", 9.75, 1.0, 3.3, 0.45,
         size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_rect(sl, 9.75, 1.45, 3.1, 0.04, WHITE)

facts = [
    ("1,000",  "country-year\nobservations"),
    ("15",     "countries\ncovered"),
    ("5",      "analytical\nlayers"),
    ("8",      "hypothesis\ntests"),
]
for i, (num, lbl) in enumerate(facts):
    yy = 1.65 + i * 1.35
    add_text(sl, num, 9.75, yy, 3.3, 0.6,
             size=36, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    add_text(sl, lbl, 9.75, yy + 0.55, 3.3, 0.55,
             size=12, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

# Bottom date
add_text(sl, "June 2026", 0.5, 6.85, 4.0, 0.4, size=12, color=MID_GREY)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 – RESEARCH BACKGROUND & RELATED WORK
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
header_bar(sl, "Research Background & Related Work",
           "What prior work frames this study?")
footer(sl)

# Left column: motivation
bullet_box(sl, 0.3, 1.5, 5.9, 2.7, "Why This Study Matters", [
    "International climate commitments (Paris Agreement) require measurable decarbonization evidence",
    "Raw climate datasets are inconsistent — country names vary, ranges are unchecked",
    "A single statistical test on dirty data is fragile; a layered pipeline is robust",
    "Reproducible methodology needed so findings can be independently verified",
])

# Right column: related work
bullet_box(sl, 6.5, 1.5, 6.5, 2.7, "Key Related Works", [
    "Ekwurzel et al. (2017) — CO₂ attribution to carbon producers; motivates data provenance",
    "Mengel et al. (2018) — committed sea level rise under Paris Agreement targets",
    "Chen et al. (2023) — data-driven pathway analysis for climate forecasting",
    "IPCC AR6 (2023) — global scientific consensus on warming trajectories",
])

# Bottom: research gap
add_rect(sl, 0.3, 4.4, 12.7, 2.55, DARK_BLUE)
add_text(sl, "Research Gap & Contribution", 0.55, 4.5, 7.0, 0.45,
         size=15, bold=True, color=GREEN)
add_text(sl,
    "Existing work focuses on individual techniques (attribution models, physical simulations, single regression). "
    "This study contributes a reproducible five-layer pipeline — normalization → EDA → segmentation → "
    "hypothesis testing → unsupervised clustering — applied to a 1,000-row observational dataset, "
    "demonstrating how multiple independent methods collectively build a statistically defensible conclusion.",
    0.55, 4.95, 12.3, 1.75, size=13, color=WHITE, wrap=True)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 – DATASET & DATA PIPELINE
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
header_bar(sl, "Dataset & Data Pipeline",
           "From flat CSV to verified relational schema")
footer(sl)

# Dataset stats strip
for i, (val, lbl) in enumerate([
        ("1,000", "Observations"), ("15", "Countries"),
        ("10", "Variables"), ("0", "Missing Values"), ("0", "FK Violations")]):
    xx = 0.3 + i * 2.55
    add_rect(sl, xx, 1.5, 2.3, 1.15, MID_BLUE)
    add_text(sl, val, xx, 1.58, 2.3, 0.62,
             size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, lbl, xx, 2.15, 2.3, 0.38,
             size=12, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

# Pipeline flow
add_text(sl, "Data Pipeline", 0.3, 2.85, 4.0, 0.42,
         size=15, bold=True, color=DARK_BLUE)

steps = [
    ("1", "Raw CSV\nclimate_change_dataset.csv", MID_BLUE),
    ("2", "SQLite 3NF Schema\ncountries + climate_data tables", DARK_BLUE),
    ("3", "Constraint Checks\nUNIQUE, FK, CHECK constraints", MID_BLUE),
    ("4", "Verification\nverify.sql — zero violations", DARK_BLUE),
    ("5", "Analysis-Ready\nNotebook loads clean data", GREEN),
]
for i, (num, lbl, col) in enumerate(steps):
    xx = 0.3 + i * 2.55
    add_rect(sl, xx, 3.3, 2.2, 1.1, col)
    add_text(sl, num, xx, 3.33, 0.48, 0.42,
             size=20, bold=True, color=GREEN if col != GREEN else WHITE,
             align=PP_ALIGN.CENTER)
    add_text(sl, lbl, xx + 0.45, 3.33, 1.65, 1.0,
             size=10.5, color=WHITE, wrap=True)
    if i < 4:
        add_text(sl, "▶", xx + 2.2, 3.65, 0.35, 0.45,
                 size=18, color=MID_BLUE, align=PP_ALIGN.CENTER)

# Variables box
bullet_box(sl, 0.3, 4.58, 12.7, 2.6, "Dataset Variables", [
    "Year (2000–2023)  |  Country (15 nations)  |  Avg Temperature (°C)",
    "CO₂ Emissions (Tons/Capita)  |  Sea Level Rise (mm)  |  Rainfall (mm)  |  Population",
    "Renewable Energy (%)  |  Extreme Weather Events  |  Forest Area (%)",
    "Derived: CO₂ Tier (Low/Medium/High)  |  Climate Risk Score (composite z-score)  |  Time Period bins",
])

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 – METHODOLOGY
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
header_bar(sl, "Methodology",
           "Five analytical layers — each adds a distinct form of evidence")
footer(sl)

layers = [
    ("Layer 1", "Data Integrity",
     ["3NF SQLite schema", "FK + UNIQUE constraints", "CHECK on numeric ranges", "verify.sql pre-check"],
     MID_BLUE),
    ("Layer 2", "Exploration (EDA)",
     ["Histograms per variable", "Scatter plots + OLS lines", "Pearson heatmap", "Time-series line charts"],
     DARK_BLUE),
    ("Layer 3", "Segmentation",
     ["CO₂ emissions tier (3 bins)", "Renewable energy tier", "Composite risk score (z-sum)", "Country-level means"],
     MID_BLUE),
    ("Layer 4", "Hypothesis Testing",
     ["Pearson + Spearman r", "Welch t-test (unequal var)", "One-way ANOVA", "OLS linear regression"],
     DARK_BLUE),
    ("Layer 5", "Unsupervised Validation",
     ["KMeans k=4", "Standardized features", "No labels supplied", "Cluster profiles vs. hypotheses"],
     GREEN),
]

for i, (tag, title, items, col) in enumerate(layers):
    xx = 0.2 + i * 2.63
    add_rect(sl, xx, 1.5, 2.48, 0.55, col)
    add_text(sl, tag, xx + 0.08, 1.52, 2.3, 0.28,
             size=10, bold=True, color=WHITE)
    add_text(sl, title, xx + 0.08, 1.78, 2.3, 0.24,
             size=13, bold=True, color=WHITE)
    card = add_rect(sl, xx, 2.1, 2.48, 4.45, WHITE)
    card.line.color.rgb = LIGHT_BLUE
    card.line.width = Pt(0.5)
    tb = sl.shapes.add_textbox(Inches(xx+0.12), Inches(2.2),
                                Inches(2.22), Inches(4.2))
    tf = tb.text_frame; tf.word_wrap = True
    for j, item in enumerate(items):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.space_before = Pt(6)
        run = p.add_run()
        run.text = f"•  {item}"
        run.font.size = Pt(12.5)
        run.font.color.rgb = DARK_GREY

# α note
add_text(sl, "All hypothesis tests: α = 0.05  |  Both Pearson and Spearman run for every correlation hypothesis",
         0.3, 6.78, 12.7, 0.35, size=11, italic=True, color=MID_GREY)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 – RESULTS: EDA & CORRELATION
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
header_bar(sl, "Results — EDA & Correlation Analysis",
           "Heatmap-guided hypothesis focus confirmed by scatter regression")
footer(sl)

# Heatmap figure
FIG_HEAT = os.path.join(FIGS, "fig1_heatmap.png")
FIG_SCAT = os.path.join(FIGS, "fig2_scatter.png")

img(sl, FIG_HEAT, 0.3, 1.45, 5.8)
add_text(sl, "Fig. 1 — Pearson Correlation Matrix", 0.3, 5.42, 5.8, 0.38,
         size=10, italic=True, color=MID_GREY, align=PP_ALIGN.CENTER)

img(sl, FIG_SCAT, 6.4, 1.45, 6.6)
add_text(sl, "Fig. 2 — Renewable Energy (%) vs CO₂ Emissions (Tons/Capita)",
         6.4, 5.42, 6.6, 0.38, size=10, italic=True, color=MID_GREY,
         align=PP_ALIGN.CENTER)

# Key findings strip
add_rect(sl, 0.3, 5.88, 12.7, 1.38, DARK_BLUE)
add_text(sl, "Key Findings", 0.55, 5.93, 3.5, 0.38,
         size=13, bold=True, color=GREEN)
findings = [
    "Negative cell between Renewable Energy and CO₂ in heatmap directed hypothesis focus before any formal test",
    "Scatter regression line: r = −0.09, confirming negative linear slope",
    "Upward trends visible in Avg Temperature and Sea Level Rise time-series",
    "Forest Area shows expected negative correlation with CO₂ emissions",
]
ftb = sl.shapes.add_textbox(Inches(0.55), Inches(6.32), Inches(12.3), Inches(0.9))
ftf = ftb.text_frame; ftf.word_wrap = True
for k, f in enumerate(findings):
    p = ftf.paragraphs[0] if k == 0 else ftf.add_paragraph()
    run = p.add_run()
    run.text = f"•  {f}"
    run.font.size = Pt(11.5)
    run.font.color.rgb = WHITE
    p.space_before = Pt(2)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 – RESULTS: STATISTICAL TESTS, ANOVA & CLUSTERING
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
header_bar(sl, "Results — Statistical Tests, ANOVA & Clustering",
           "Convergence of five independent methods")
footer(sl)

FIG_PVAL = os.path.join(FIGS, "fig3_pvalues.png")
FIG_CO2  = os.path.join(FIGS, "fig4_co2_country.png")

img(sl, FIG_PVAL, 0.3, 1.45, 6.5)
add_text(sl, "Fig. 3 — Hypothesis Test p-values (log scale, α = 0.45)",
         0.3, 5.05, 6.5, 0.38, size=10, italic=True, color=MID_GREY,
         align=PP_ALIGN.CENTER)

img(sl, FIG_CO2, 6.8, 1.45, 6.2)
add_text(sl, "Fig. 4 — Mean CO₂ by Country + Top/Bottom 3 Emitter Trends",
         6.8, 5.05, 6.2, 0.38, size=10, italic=True, color=MID_GREY,
         align=PP_ALIGN.CENTER)

# Results summary table
add_rect(sl, 0.3, 5.5, 12.7, 1.72, WHITE)
headers = ["Test", "Hypotheses", "Result", "Decision"]
col_xs  = [0.42, 2.2, 7.5, 10.2]
col_ws  = [1.7,  5.2, 2.6,  2.9]

add_rect(sl, 0.3, 5.5, 12.7, 0.38, DARK_BLUE)
for hdr, cx in zip(headers, col_xs):
    add_text(sl, hdr, cx, 5.53, 2.8, 0.3,
             size=11, bold=True, color=WHITE)

rows = [
    ("Welch t-test",   "High vs Low renewable CO₂ means",           "Groups differ significantly",       "✅ Reject H0"),
    ("ANOVA",          "CO₂ emissions across 15 countries",          "F significant, p < 0.05",           "✅ Reject H0"),
    ("Linear Reg.",    "Temperature & sea level trend over years",       "Positive slope, p < 0.05 both",     "✅ Reject H0"),
    ("KMeans k=4",     "Cluster 1: high renewables, low CO₂",       "Unsupervised corroboration",        "✅ Confirmed"),
]
for r, (test, hyp, res, dec) in enumerate(rows):
    yy = 5.9 + r * 0.32
    bg_col = OFF_WHITE if r % 2 == 0 else WHITE
    add_rect(sl, 0.3, yy, 12.7, 0.32, bg_col)
    for txt, cx in zip([test, hyp, res, dec], col_xs):
        col = GREEN if "✅" in txt else DARK_GREY
        add_text(sl, txt, cx, yy + 0.04, 2.8, 0.26, size=10, color=col)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 – CONCLUSION & FUTURE WORK
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
header_bar(sl, "Conclusion & Future Work",
           "Key takeaways and next steps")
footer(sl)

# Main conclusion box
add_rect(sl, 0.3, 1.5, 12.7, 1.62, DARK_BLUE)
add_text(sl, "Main Conclusion", 0.55, 1.55, 5.0, 0.42,
         size=14, bold=True, color=GREEN)
add_text(sl,
    "A five-layer analytical pipeline — normalization → EDA → segmentation → hypothesis testing "
    "→ KMeans clustering — applied to 1,000 country-year observations produced statistically robust evidence "
    "that higher renewable energy adoption is associated with significantly lower CO₂ emissions per capita. "
    "No single technique produced this finding alone; the conclusion emerges from the convergence of five "
    "independent methods pointing in the same direction.",
    0.55, 1.97, 12.1, 1.1, size=13, color=WHITE, wrap=True)

# Left: Lessons learned
bullet_box(sl, 0.3, 3.3, 5.9, 3.05, "Lessons Learned", [
    "Schema enforcement is an analytical decision — it determines whether downstream statistics are trustworthy",
    "EDA should precede hypothesis selection: the heatmap guided focus before tests were specified",
    "Running both Pearson and Spearman on the same hypothesis is a low-cost robustness check",
    "Unsupervised methods (KMeans) serve as a bias-free independent check on hypothesis-driven findings",
    "Country-year aggregation masks sub-national variation — a key limitation to address in future work",
])

# Right: Future work
bullet_box(sl, 6.5, 3.3, 6.5, 3.05, "Future Research Directions", [
    "Causal inference — Granger causality or instrumental variables to move beyond correlation",
    "Finer temporal resolution — monthly data to enable seasonal decomposition",
    "Additional confounders — join GDP per capita and energy policy indices",
    "Forecasting — LSTM or Prophet models for temperature and sea level projections",
    "Optimal k validation — elbow plot and silhouette score to choose KMeans k empirically",
])

# ─── Save ─────────────────────────────────────────────────────────────────────
prs.save(OUTPUT)
print(f"Presentation saved: {OUTPUT}")
